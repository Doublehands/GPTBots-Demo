"""商业调查：按文件顺序解析 Excel 行、PPTX 幻灯片、PDF 页面文本。"""

from __future__ import annotations

import io
import re
from typing import Any

from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader


def _excel_rows_from_xlsx(data: bytes, *, skip_first_row: bool) -> list[dict[str, Any]]:
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    items: list[dict[str, Any]] = []
    try:
        ws = wb.active
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if skip_first_row and i == 0:
                continue
            cells = ["" if c is None else str(c).strip() for c in row]
            if not any(cells):
                continue
            line = "\t".join(cells)
            items.append(
                {
                    "kind": "excel_row",
                    "label": f"表格 第 {i + 1} 行",
                    "text": line,
                    "meta": {"row_index": i + 1},
                }
            )
    finally:
        wb.close()
    return items


def _pptx_slides(data: bytes, filename: str) -> list[dict[str, Any]]:
    prs = Presentation(io.BytesIO(data))
    items: list[dict[str, Any]] = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for p in shape.text_frame.paragraphs:
                t = "".join(run.text for run in p.runs).strip()
                if t:
                    parts.append(t)
        text = "\n".join(parts).strip()
        if not text:
            text = f"（第 {idx} 页无文本）"
        items.append(
            {
                "kind": "ppt_slide",
                "label": f"{filename} · 第 {idx} 页",
                "text": text,
                "meta": {"slide_index": idx, "file": filename},
            }
        )
    return items


def _pdf_pages(data: bytes, filename: str) -> list[dict[str, Any]]:
    reader = PdfReader(io.BytesIO(data))
    items: list[dict[str, Any]] = []
    for i, page in enumerate(reader.pages, start=1):
        raw = page.extract_text() or ""
        text = re.sub(r"\s+", " ", raw).strip()
        if not text:
            text = f"（第 {i} 页未提取到文本，可换用 OCR 方案）"
        items.append(
            {
                "kind": "pdf_page",
                "label": f"{filename} · 第 {i} 页",
                "text": text,
                "meta": {"page_index": i, "file": filename},
            }
        )
    return items


def build_queue_from_files(
    ordered_files: list[tuple[str, bytes]],
    *,
    skip_first_excel_row: bool,
) -> tuple[list[dict[str, Any]], list[str]]:
    """
    按 ordered_files 顺序依次展开：每个 Excel 逐行、每个 PPTX 逐页幻灯片、每个 PDF 逐页。
    返回 (items, warnings)。
    """
    items: list[dict[str, Any]] = []
    warnings: list[str] = []

    for filename, data in ordered_files:
        lower = filename.lower()
        if lower.endswith(".xlsx") or lower.endswith(".xlsm"):
            rows = _excel_rows_from_xlsx(data, skip_first_row=skip_first_excel_row)
            for r in rows:
                r["meta"]["file"] = filename
                r["label"] = f"{filename} · {r['label']}"
            items.extend(rows)
        elif lower.endswith(".pptx"):
            items.extend(_pptx_slides(data, filename))
        elif lower.endswith(".pdf"):
            items.extend(_pdf_pages(data, filename))
        elif lower.endswith(".ppt"):
            warnings.append(f"已跳过 {filename}：旧版 .ppt 不支持，请另存为 .pptx")
        elif lower.endswith(".xls"):
            warnings.append(f"已跳过 {filename}：.xls 请另存为 .xlsx 后再上传")
        else:
            warnings.append(f"已跳过 {filename}：仅支持 .xlsx / .pptx / .pdf")

    return items, warnings
